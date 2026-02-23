#!/usr/bin/env python3
"""
generate_quiz_pptx.py
---------------------
Generates an interactive quiz-style PowerPoint lecture for psychiatry residents
covering eating disorders board exam content.

Usage:
    python generate_quiz_pptx.py
    python generate_quiz_pptx.py --category "Pharmacotherapy"
    python generate_quiz_pptx.py --output my_quiz.pptx
    python generate_quiz_pptx.py --format jeopardy
    python generate_quiz_pptx.py --format audience_response

Requirements:
    pip install python-pptx
"""

import argparse
import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme colours (deep navy + gold accent, professional medical feel)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0D, 0x2B, 0x55)       # slide background / title bars
GOLD = RGBColor(0xC9, 0xA0, 0x2C)       # accent / highlights
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF8)
GREEN = RGBColor(0x1A, 0x7A, 0x4A)      # correct-answer reveal
RED = RGBColor(0xB3, 0x1B, 0x1B)        # distractor / wrong

QUIZ_BANK_PATH = Path(__file__).parent / "questions" / "quiz_bank.json"


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def _load_quiz_bank(path: Path) -> dict:
    with open(path, "r", encoding="utf-8") as fh:
        return json.load(fh)


def _set_slide_background(slide, colour: RGBColor):
    """Fill the slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int = 18, bold: bool = False,
                  color: RGBColor = WHITE, align=PP_ALIGN.LEFT,
                  wrap: bool = True) -> object:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _add_rect(slide, left, top, width, height,
              fill_color: RGBColor, line_color: RGBColor = None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, NAVY)

    W = prs.slide_width
    H = prs.slide_height

    # Gold accent bar at top
    _add_rect(slide, 0, 0, W, Inches(0.12), GOLD)

    # Main title
    _add_text_box(slide, title,
                  Inches(0.6), Inches(1.8), W - Inches(1.2), Inches(1.8),
                  font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    _add_text_box(slide, subtitle,
                  Inches(0.6), Inches(3.8), W - Inches(1.2), Inches(1.2),
                  font_size=24, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

    # Gold accent bar at bottom
    _add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), GOLD)

    return slide


def add_section_divider(prs: Presentation, section_name: str):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, NAVY)

    W = prs.slide_width
    H = prs.slide_height

    _add_rect(slide, Inches(0.5), H / 2 - Inches(0.06),
              W - Inches(1.0), Inches(0.06), GOLD)

    _add_text_box(slide, section_name,
                  Inches(0.5), H / 2 - Inches(1.0),
                  W - Inches(1.0), Inches(0.9),
                  font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    return slide


def add_instructions_slide(prs: Presentation, format_name: str):
    """Add a brief instructions slide for the chosen quiz format."""
    instructions = {
        "standard": (
            "How to Use These Slides",
            [
                "Each question slide shows the stem and answer choices.",
                "Allow the audience to discuss or vote before clicking to advance.",
                "The NEXT slide reveals the correct answer with explanation.",
                "Tip: Use PowerPoint Presenter View to see notes while presenting.",
            ]
        ),
        "jeopardy": (
            "Jeopardy-Style Game Instructions",
            [
                "Categories are shown on the game board slide.",
                "Teams take turns choosing a category and point value.",
                "Presenter advances to the question, team answers, then advance to reveal.",
                "Keep score on a whiteboard or use the score tracker at the end.",
            ]
        ),
        "audience_response": (
            "Audience Response System Instructions",
            [
                "Use Poll Everywhere, Mentimeter, or similar tool for live voting.",
                "QR codes or URLs can be embedded in each question slide.",
                "Display results before revealing the correct answer for discussion.",
                "Export poll results to track resident performance over time.",
            ]
        ),
    }
    title, bullets = instructions.get(format_name, instructions["standard"])

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, LIGHT_GREY)

    W = prs.slide_width
    H = prs.slide_height

    _add_rect(slide, 0, 0, W, Inches(1.1), NAVY)
    _add_text_box(slide, title,
                  Inches(0.4), Inches(0.15), W - Inches(0.8), Inches(0.8),
                  font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    for i, bullet in enumerate(bullets):
        _add_text_box(slide, f"▸  {bullet}",
                      Inches(0.7), Inches(1.4) + i * Inches(0.85),
                      W - Inches(1.2), Inches(0.75),
                      font_size=18, color=NAVY, align=PP_ALIGN.LEFT)

    return slide


def add_multiple_choice_question(prs: Presentation, q: dict, q_num: int):
    """
    Creates TWO slides per question:
      1. Question + answer choices (participants think / respond)
      2. Correct-answer reveal + explanation
    """
    slide_layout = prs.slide_layouts[6]
    W = prs.slide_width
    H = prs.slide_height

    # ---- Slide 1: Question ------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, LIGHT_GREY)

    # Header bar
    _add_rect(slide, 0, 0, W, Inches(1.05), NAVY)
    _add_text_box(slide,
                  f"Q{q_num}  |  {q.get('board_topic', '')}",
                  Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.45),
                  font_size=14, bold=False, color=GOLD, align=PP_ALIGN.LEFT)
    difficulty_label = q.get("difficulty", "").upper()
    _add_text_box(slide, f"Difficulty: {difficulty_label}",
                  W - Inches(2.0), Inches(0.1), Inches(1.8), Inches(0.45),
                  font_size=13, bold=False, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)

    # Question stem
    _add_text_box(slide, q["question"],
                  Inches(0.4), Inches(1.15), W - Inches(0.8), Inches(1.6),
                  font_size=20, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    # Answer choices
    choice_colors = {
        "A": RGBColor(0x1A, 0x4A, 0x7A),
        "B": RGBColor(0x1A, 0x4A, 0x7A),
        "C": RGBColor(0x1A, 0x4A, 0x7A),
        "D": RGBColor(0x1A, 0x4A, 0x7A),
    }
    choices = q.get("choices", {})
    positions = [
        (Inches(0.4), Inches(2.9)),
        (Inches(0.4), Inches(3.75)),
        (Inches(0.4), Inches(4.6)),
        (Inches(0.4), Inches(5.45)),
    ]
    for idx, (letter, text) in enumerate(choices.items()):
        if idx >= len(positions):
            break
        lx, ly = positions[idx]
        # coloured pill
        _add_rect(slide, lx, ly, W - Inches(0.8), Inches(0.75),
                  choice_colors.get(letter, NAVY))
        _add_text_box(slide, f"{letter}.  {text}",
                      lx + Inches(0.15), ly + Inches(0.08),
                      W - Inches(1.1), Inches(0.6),
                      font_size=17, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    # Prompt at bottom
    _add_text_box(slide, "⏱  Discuss with your team, then advance to reveal the answer.",
                  Inches(0.4), H - Inches(0.55), W - Inches(0.8), Inches(0.45),
                  font_size=13, bold=False, color=NAVY, align=PP_ALIGN.CENTER)

    # ---- Slide 2: Answer Reveal -------------------------------------------
    reveal_slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(reveal_slide, LIGHT_GREY)

    # Header bar
    _add_rect(reveal_slide, 0, 0, W, Inches(1.05), GREEN)
    _add_text_box(reveal_slide,
                  f"Q{q_num}  ANSWER REVEAL  |  {q.get('board_topic', '')}",
                  Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.45),
                  font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Restate question (smaller)
    _add_text_box(reveal_slide, q["question"],
                  Inches(0.4), Inches(1.15), W - Inches(0.8), Inches(1.1),
                  font_size=16, bold=False, color=NAVY, align=PP_ALIGN.LEFT)

    # Re-render choices; highlight correct, grey out others
    correct = q.get("answer", "")
    for idx, (letter, text) in enumerate(choices.items()):
        if idx >= len(positions):
            break
        lx, ly = positions[idx]
        if letter == correct:
            fill = GREEN
            icon = "✓"
            fw = True
        else:
            fill = RGBColor(0xCC, 0xCC, 0xCC)
            icon = " "
            fw = False
        _add_rect(reveal_slide, lx, ly, W - Inches(0.8), Inches(0.72), fill)
        _add_text_box(reveal_slide, f"{icon} {letter}.  {text}",
                      lx + Inches(0.12), ly + Inches(0.06),
                      W - Inches(1.1), Inches(0.6),
                      font_size=17, bold=fw, color=WHITE, align=PP_ALIGN.LEFT)

    # Explanation box
    _add_rect(reveal_slide, Inches(0.35), H - Inches(1.85),
              W - Inches(0.7), Inches(1.65), NAVY)
    _add_text_box(reveal_slide,
                  f"📚  {q.get('explanation', '')}",
                  Inches(0.5), H - Inches(1.8),
                  W - Inches(1.0), Inches(1.55),
                  font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    return slide, reveal_slide


def add_true_false_question(prs: Presentation, q: dict, q_num: int):
    """Creates TWO slides: question + answer reveal for True/False format."""
    slide_layout = prs.slide_layouts[6]
    W = prs.slide_width
    H = prs.slide_height

    # ---- Slide 1: Question ------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, LIGHT_GREY)

    _add_rect(slide, 0, 0, W, Inches(1.05), NAVY)
    _add_text_box(slide, f"Q{q_num}  |  True or False?  |  {q.get('board_topic', '')}",
                  Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.45),
                  font_size=14, bold=False, color=GOLD)

    _add_text_box(slide, q["question"],
                  Inches(0.4), Inches(1.2), W - Inches(0.8), Inches(2.0),
                  font_size=22, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    # TRUE / FALSE buttons
    _add_rect(slide, Inches(0.8), Inches(3.5), Inches(3.6), Inches(1.1), GREEN)
    _add_text_box(slide, "TRUE", Inches(0.8), Inches(3.6), Inches(3.6), Inches(0.9),
                  font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_rect(slide, Inches(5.5), Inches(3.5), Inches(3.6), Inches(1.1), RED)
    _add_text_box(slide, "FALSE", Inches(5.5), Inches(3.6), Inches(3.6), Inches(0.9),
                  font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_text_box(slide, "⏱  Vote now, then advance to the answer.",
                  Inches(0.4), H - Inches(0.55), W - Inches(0.8), Inches(0.45),
                  font_size=13, color=NAVY, align=PP_ALIGN.CENTER)

    # ---- Slide 2: Answer Reveal -------------------------------------------
    reveal = prs.slides.add_slide(slide_layout)
    _set_slide_background(reveal, LIGHT_GREY)

    answer_is_true = q.get("answer", "").strip().lower() == "true"
    header_color = GREEN if answer_is_true else RED
    answer_text = "✓  TRUE" if answer_is_true else "✗  FALSE"

    _add_rect(reveal, 0, 0, W, Inches(1.05), header_color)
    _add_text_box(reveal, f"Q{q_num}  ANSWER REVEAL  |  {q.get('board_topic', '')}",
                  Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.45),
                  font_size=14, bold=True, color=WHITE)

    _add_text_box(reveal, q["question"],
                  Inches(0.4), Inches(1.2), W - Inches(0.8), Inches(1.5),
                  font_size=18, bold=False, color=NAVY)

    _add_text_box(reveal, answer_text,
                  Inches(0.4), Inches(2.85), W - Inches(0.8), Inches(0.75),
                  font_size=36, bold=True, color=header_color, align=PP_ALIGN.CENTER)

    _add_rect(reveal, Inches(0.35), H - Inches(2.1),
              W - Inches(0.7), Inches(1.9), NAVY)
    _add_text_box(reveal,
                  f"📚  {q.get('explanation', '')}",
                  Inches(0.5), H - Inches(2.05),
                  W - Inches(1.0), Inches(1.8),
                  font_size=14, color=WHITE)

    return slide, reveal


def add_case_vignette_question(prs: Presentation, q: dict, q_num: int):
    """
    Case vignette: three slides
      1. Clinical stem (read / discuss)
      2. Question + choices
      3. Answer reveal + explanation
    """
    slide_layout = prs.slide_layouts[6]
    W = prs.slide_width
    H = prs.slide_height

    # ---- Slide 1: Clinical Stem -------------------------------------------
    stem_slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(stem_slide, LIGHT_GREY)

    _add_rect(stem_slide, 0, 0, W, Inches(1.05), NAVY)
    _add_text_box(stem_slide,
                  f"CASE VIGNETTE  Q{q_num}  |  {q.get('board_topic', '')}",
                  Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.45),
                  font_size=14, bold=True, color=GOLD)

    # Clinical stem text box with light background
    _add_rect(stem_slide, Inches(0.3), Inches(1.15),
              W - Inches(0.6), H - Inches(1.75), NAVY)
    _add_text_box(stem_slide, q.get("clinical_stem", ""),
                  Inches(0.5), Inches(1.25),
                  W - Inches(1.0), H - Inches(2.0),
                  font_size=17, color=WHITE, align=PP_ALIGN.LEFT)

    _add_text_box(stem_slide,
                  "Read the case, then advance to the question.",
                  Inches(0.4), H - Inches(0.5),
                  W - Inches(0.8), Inches(0.4),
                  font_size=13, color=NAVY, align=PP_ALIGN.CENTER)

    # ---- Slide 2 & 3: reuse MC logic (without clinical_stem on q slide) ----
    q_slide, reveal_slide = add_multiple_choice_question(prs, q, q_num)

    return stem_slide, q_slide, reveal_slide


def add_jeopardy_board(prs: Presentation, categories: list):
    """
    Creates a simple Jeopardy game-board slide with category columns
    and point rows (100, 200, 300, 400, 500).
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, NAVY)

    W = prs.slide_width
    H = prs.slide_height

    _add_text_box(slide, "EATING DISORDERS  JEOPARDY",
                  Inches(0.3), Inches(0.05), W - Inches(0.6), Inches(0.65),
                  font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    points = [100, 200, 300, 400, 500]
    n_cats = min(len(categories), 5)
    col_w = (W - Inches(0.4)) / n_cats
    row_h = (H - Inches(0.85)) / (len(points) + 1)

    # Category headers
    for ci, cat in enumerate(categories[:n_cats]):
        lx = Inches(0.2) + ci * col_w
        _add_rect(slide, lx + Inches(0.05), Inches(0.75),
                  col_w - Inches(0.1), row_h - Inches(0.08), GOLD)
        _add_text_box(slide, cat,
                      lx + Inches(0.08), Inches(0.78),
                      col_w - Inches(0.16), row_h - Inches(0.14),
                      font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Point cells
    for ri, pts in enumerate(points):
        row_top = Inches(0.75) + (ri + 1) * row_h
        for ci in range(n_cats):
            lx = Inches(0.2) + ci * col_w
            _add_rect(slide, lx + Inches(0.05), row_top + Inches(0.04),
                      col_w - Inches(0.1), row_h - Inches(0.12), GOLD)
            _add_text_box(slide, f"${pts}",
                          lx + Inches(0.08), row_top + Inches(0.07),
                          col_w - Inches(0.16), row_h - Inches(0.18),
                          font_size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    return slide


def add_score_tracker(prs: Presentation, n_teams: int = 4):
    """Adds a simple score-tracker slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, NAVY)

    W = prs.slide_width
    H = prs.slide_height

    _add_text_box(slide, "SCORE TRACKER",
                  Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.7),
                  font_size=34, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    team_names = [f"Team {i + 1}" for i in range(n_teams)]
    col_w = (W - Inches(0.6)) / n_teams

    for ti, name in enumerate(team_names):
        lx = Inches(0.3) + ti * col_w
        _add_rect(slide, lx + Inches(0.1), Inches(0.95),
                  col_w - Inches(0.2), Inches(0.65), GOLD)
        _add_text_box(slide, name,
                      lx + Inches(0.12), Inches(1.0),
                      col_w - Inches(0.24), Inches(0.55),
                      font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # Score area
        _add_rect(slide, lx + Inches(0.1), Inches(1.65),
                  col_w - Inches(0.2), H - Inches(2.2),
                  RGBColor(0x1A, 0x3A, 0x6A))
        _add_text_box(slide, "0",
                      lx + Inches(0.12), Inches(1.9),
                      col_w - Inches(0.24), Inches(1.2),
                      font_size=48, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    return slide


def add_key_facts_slide(prs: Presentation, facts: list, title: str = "High-Yield Board Facts"):
    """Summary slide with bullet-point key facts."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, LIGHT_GREY)

    W = prs.slide_width
    H = prs.slide_height

    _add_rect(slide, 0, 0, W, Inches(1.05), NAVY)
    _add_text_box(slide, title,
                  Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.8),
                  font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, fact in enumerate(facts[:7]):
        _add_text_box(slide, f"★  {fact}",
                      Inches(0.5), Inches(1.2) + i * Inches(0.72),
                      W - Inches(1.0), Inches(0.65),
                      font_size=16, color=NAVY, align=PP_ALIGN.LEFT)

    return slide


# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------

def build_presentation(quiz_bank: dict, category_filter: str = None,
                        fmt: str = "standard", output_path: str = None) -> str:
    """
    Build the full quiz PowerPoint presentation.

    Parameters
    ----------
    quiz_bank       : loaded JSON quiz bank
    category_filter : optional category name substring to limit questions
    fmt             : 'standard' | 'jeopardy' | 'audience_response'
    output_path     : output .pptx file path (default: eating_disorders_quiz.pptx)
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    metadata = quiz_bank.get("metadata", {})
    output_path = output_path or "eating_disorders_quiz.pptx"

    # ---- Title slide -------------------------------------------------------
    add_title_slide(
        prs,
        metadata.get("title", "Eating Disorders Quiz"),
        f"Psychiatry Residency Board Review\n{metadata.get('last_updated', '')}",
    )

    # ---- Instructions slide ------------------------------------------------
    add_instructions_slide(prs, fmt)

    # ---- Jeopardy game board (jeopardy format only) -----------------------
    categories = quiz_bank.get("categories", [])
    if category_filter:
        categories = [
            c for c in categories
            if category_filter.lower() in c["name"].lower()
        ]

    if fmt == "jeopardy":
        cat_names = [c["name"] for c in categories][:5]
        add_jeopardy_board(prs, cat_names)
        add_score_tracker(prs)

    # ---- Question slides ---------------------------------------------------
    q_num = 0

    for category in categories:
        add_section_divider(prs, category["name"])

        for q in category.get("questions", []):
            q_num += 1
            q_type = q.get("type", "multiple_choice")

            if q_type == "multiple_choice":
                add_multiple_choice_question(prs, q, q_num)
            elif q_type == "true_false":
                add_true_false_question(prs, q, q_num)
            elif q_type == "case_vignette":
                add_case_vignette_question(prs, q, q_num)
            else:
                add_multiple_choice_question(prs, q, q_num)

    # ---- High-yield facts summary -----------------------------------------
    high_yield_facts = [
        "AN has the highest mortality rate of ANY psychiatric disorder (SMR ~5-10x).",
        "DSM-5 REMOVED amenorrhea as a required criterion for AN.",
        "DSM-5 reduced BN binge/purge frequency threshold: once/week (was twice/week).",
        "Fluoxetine 60 mg/day = only FDA-approved med for Bulimia Nervosa.",
        "Lisdexamfetamine (Vyvanse) = only FDA-approved med for Binge Eating Disorder.",
        "BUPROPION is CONTRAINDICATED in BN (↑ seizure risk).",
        "Refeeding syndrome: hypophosphatemia → cardiac/respiratory failure.",
        "Russell's Sign = dorsal hand calluses from self-induced vomiting (BN).",
        "FBT (Maudsley) = first-line therapy for adolescents with AN.",
        "CBT-E = first-line psychotherapy for BN and BED in adults.",
    ]
    add_key_facts_slide(prs, high_yield_facts[:7], "High-Yield Board Facts – Part 1")
    add_key_facts_slide(prs, high_yield_facts[7:], "High-Yield Board Facts – Part 2")

    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description=(
            "Generate an interactive quiz PowerPoint for eating disorder "
            "psychiatry board review."
        )
    )
    parser.add_argument(
        "--category",
        type=str,
        default=None,
        help="Filter questions to a specific category (partial match, case-insensitive). "
             "E.g. --category 'Pharmacotherapy'",
    )
    parser.add_argument(
        "--output",
        type=str,
        default="eating_disorders_quiz.pptx",
        help="Output PowerPoint filename (default: eating_disorders_quiz.pptx)",
    )
    parser.add_argument(
        "--format",
        dest="fmt",
        type=str,
        default="standard",
        choices=["standard", "jeopardy", "audience_response"],
        help=(
            "Quiz format: 'standard' (Q→answer pairs), "
            "'jeopardy' (game board + scoring), "
            "'audience_response' (live-vote prompts). "
            "Default: standard"
        ),
    )
    parser.add_argument(
        "--bank",
        type=str,
        default=str(QUIZ_BANK_PATH),
        help="Path to quiz bank JSON file (default: questions/quiz_bank.json)",
    )

    args = parser.parse_args()

    bank_path = Path(args.bank)
    if not bank_path.exists():
        print(f"ERROR: Quiz bank not found at {bank_path}", file=sys.stderr)
        sys.exit(1)

    print(f"Loading quiz bank from: {bank_path}")
    quiz_bank = _load_quiz_bank(bank_path)

    print(f"Building presentation  [format={args.fmt}] …")
    output = build_presentation(
        quiz_bank,
        category_filter=args.category,
        fmt=args.fmt,
        output_path=args.output,
    )

    print(f"✅  Saved: {output}")
    n_slides = len(__import__("pptx").Presentation(output).slides)
    print(f"   Total slides: {n_slides}")


if __name__ == "__main__":
    main()
